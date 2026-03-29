#!/usr/bin/env python3
"""
Generate PowerPoint presentation from technology comparison data.

Usage:
    python create_comparison_presentation.py <comparison_data.json> <output.pptx>

Requires: python-pptx
Install: pip3 install python-pptx

Input JSON format:
{
    "title": "Technology Comparison",
    "date": "2026-03-27",
    "technologies": ["Tech A", "Tech B", "Tech C"],
    "executive_summary": "...",
    "scenarios": [...],
    "dimensions": [...],
    "recommendations": {...}
}
"""

import json
import sys
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# Brand colors
PRIMARY_COLOR = RGBColor(238, 0, 0)  # Red Hat Red
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)


def create_title_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    """Create title slide with disclaimer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = data.get('title', 'Technology Comparison')
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_text = f"{', '.join(data.get('technologies', []))}"
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = data.get('date', '')
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = DARK_GRAY
    date_para.alignment = PP_ALIGN.CENTER

    # Disclaimer
    disclaimer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
    disclaimer_frame = disclaimer_box.text_frame
    disclaimer_frame.text = "⚠️ This report is AI generated with minimal human input and validation."
    disclaimer_para = disclaimer_frame.paragraphs[0]
    disclaimer_para.font.size = Pt(14)
    disclaimer_para.font.italic = True
    disclaimer_para.font.color.rgb = DARK_GRAY
    disclaimer_para.alignment = PP_ALIGN.CENTER


def create_header(slide, title: str) -> None:
    """Add header to slide."""
    # Header background
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR

    # Header text
    header_text = header.text_frame
    header_text.text = title
    header_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = header_text.paragraphs[0]
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.LEFT
    para.level = 0
    # Add left margin
    header_text.margin_left = Inches(0.3)


def add_bullet_points(text_frame, items: List[str], font_size: int = 18) -> None:
    """Add bullet points to text frame."""
    text_frame.clear()

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)


def create_scenario_slides(prs: Presentation, scenarios: List[Dict[str, Any]]) -> None:
    """Create slides for scenario analysis."""

    # Group scenarios into slides (max 5 per slide)
    scenario_groups = [scenarios[i:i+5] for i in range(0, len(scenarios), 5)]

    for idx, group in enumerate(scenario_groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        title = f"Scenario Analysis ({idx+1}/{len(scenario_groups)})" if len(scenario_groups) > 1 else "Scenario Analysis"
        create_header(slide, title)

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.2))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        bullets = []
        for scenario in group:
            scenario_name = scenario.get('scenario', scenario.get('name', 'Unknown'))
            priority = scenario.get('priority', '')
            best_fit = scenario.get('best_fit', 'N/A')

            # Priority emoji
            priority_emoji = ""
            if priority == "CRITICAL":
                priority_emoji = "🔴"
            elif priority == "IMPORTANT":
                priority_emoji = "🟡"
            elif priority == "NICE-TO-HAVE":
                priority_emoji = "⚪"

            bullets.append(f"{priority_emoji} {scenario_name}: {best_fit}")

        add_bullet_points(text_frame, bullets, font_size=20)


def create_dimension_slides(prs: Presentation, dimensions: List[Dict[str, Any]], technologies: List[str]) -> None:
    """Create slides for dimension comparisons."""

    # Group dimensions into slides (max 4 per slide)
    dimension_groups = [dimensions[i:i+4] for i in range(0, len(dimensions), 4)]

    for idx, group in enumerate(dimension_groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        title = f"Comparison Dimensions ({idx+1}/{len(dimension_groups)})" if len(dimension_groups) > 1 else "Comparison Dimensions"
        create_header(slide, title)

        # Content area - create a table
        rows = len(group) + 1  # +1 for header
        cols = len(technologies) + 1  # +1 for dimension name

        # Table positioning
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row
        table.cell(0, 0).text = "Dimension"
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(14)
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = LIGHT_GRAY

        for col_idx, tech in enumerate(technologies, start=1):
            cell = table.cell(0, col_idx)
            cell.text = tech
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

        # Dimension rows
        for row_idx, dimension in enumerate(group, start=1):
            dim_name = dimension.get('name', 'Unknown')
            table.cell(row_idx, 0).text = dim_name
            table.cell(row_idx, 0).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(row_idx, 0).text_frame.paragraphs[0].font.bold = True

            comparisons = dimension.get('comparisons', {})
            for col_idx, tech in enumerate(technologies, start=1):
                comp = comparisons.get(tech, {})
                score = comp.get('score', 'N/A')

                # Convert to emoji
                if score in ["Strong", "Excellent", "Best"]:
                    display = "✅"
                elif score in ["Good", "Solid"]:
                    display = "✅"
                elif score in ["Moderate", "Fair", "Acceptable"]:
                    display = "⚠️"
                elif score in ["Weak", "Limited", "Poor"]:
                    display = "❌"
                else:
                    display = score

                cell = table.cell(row_idx, col_idx)
                cell.text = display
                cell.text_frame.paragraphs[0].font.size = Pt(16)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_recommendation_slide(prs: Presentation, recommendation: Dict[str, Any]) -> None:
    """Create recommendation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    create_header(slide, "Recommendations")

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    # Primary recommendation
    p1 = text_frame.paragraphs[0]
    p1.text = f"Primary Recommendation: {recommendation.get('primary', 'N/A')}"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_COLOR
    p1.space_after = Pt(12)

    # Rationale
    p2 = text_frame.add_paragraph()
    p2.text = recommendation.get('rationale', '')
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(12)

    # Alternatives
    if 'alternatives' in recommendation and recommendation['alternatives']:
        p3 = text_frame.add_paragraph()
        p3.text = "Alternatives to Consider:"
        p3.font.size = Pt(18)
        p3.font.bold = True
        p3.font.color.rgb = DARK_GRAY
        p3.space_after = Pt(6)

        for alt in recommendation['alternatives'][:2]:  # Top 2
            p_alt = text_frame.add_paragraph()
            p_alt.text = f"• {alt.get('technology', 'Unknown')}"
            p_alt.font.size = Pt(16)
            p_alt.font.color.rgb = DARK_GRAY
            p_alt.level = 1


def generate_presentation(data: Dict[str, Any]) -> Presentation:
    """Generate complete PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(prs, data)

    # Slide 2+: Scenario Analysis
    if 'scenario_analysis' in data or 'scenarios' in data:
        scenarios = data.get('scenario_analysis', data.get('scenarios', []))
        if scenarios:
            create_scenario_slides(prs, scenarios)

    # Slides: Dimension Comparison
    if 'dimensions' in data:
        dimensions = data.get('dimensions', [])
        technologies = data.get('technologies', [])
        if dimensions and technologies:
            create_dimension_slides(prs, dimensions, technologies)

    # Final Slide: Recommendations
    if 'overall_recommendation' in data:
        create_recommendation_slide(prs, data['overall_recommendation'])
    elif 'recommendations' in data and isinstance(data['recommendations'], dict):
        create_recommendation_slide(prs, data['recommendations'])

    return prs


def main():
    if len(sys.argv) != 3:
        print("Usage: python create_comparison_presentation.py <comparison_data.json> <output.pptx>")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]

    # Load comparison data
    try:
        with open(input_file, 'r') as f:
            data = json.load(f)
    except FileNotFoundError:
        print(f"Error: Input file '{input_file}' not found")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in input file: {e}")
        sys.exit(1)

    # Check for python-pptx
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx not installed. Run: pip3 install python-pptx")
        sys.exit(1)

    # Generate presentation
    prs = generate_presentation(data)

    # Save
    prs.save(output_file)
    print(f"✅ Presentation generated: {output_file}")


if __name__ == "__main__":
    main()
